# # main.py

# import os
# from tqdm import tqdm
# from langchain_community.document_loaders import AzureBlobStorageContainerLoader
# from agents.reader_agent import run_reader_agent
# from agents.categorizer_agent import run_categorization_agent

# # Load documents from Azure Blob Storage
# loader = AzureBlobStorageContainerLoader(
#     conn_str=os.getenv("AZURE_STORAGE_CONNECTION_STRING"),
#     container=os.getenv("AZURE_STORAGE_CONTAINER_NAME")
# )
# documents = loader.load()

# def process_documents(documents):
#     results = []
#     categories = []

#     for doc in tqdm(documents):
#         try:
#             source_path = doc.metadata.get("source")
#             file_name = os.path.basename(source_path)
#             if file_name.lower().endswith(('.pdf', '.pptx')):
#                 text = doc.page_content

#                 # Agent 1: Read and Summarize
#                 summary = run_reader_agent(text)
#                 print(f"\nSummary of {file_name}:\n{summary}")

#                 # Agent 2: Categorize
#                 category = run_categorization_agent(summary, categories)
#                 if category not in categories:
#                     categories.append(category)

#                 results.append({"filename": file_name, "category": category})
#         except Exception as e:
#             print(f"Error processing {file_name}: {e}")

#     return results, categories

# if __name__ == "__main__":
#     results, all_categories = process_documents(documents)

#     print("\nCategorization Results:")
#     for item in results:
#         print(f"{item['filename']}: {item['category']}")

#     print("\nFinal Categories Identified:")
#     for c in all_categories:
#         print("-", c)

# main.py
# import os
# import streamlit as st
# from langchain_community.document_loaders import AzureBlobStorageContainerLoader
# from agents.reader_agent import run_reader_agent
# from agents.categorizer_agent import run_categorization_agent

# # === Streamlit Setup ===
# st.set_page_config(page_title="Case Study Categorizer", layout="wide")
# st.title("📁 Case Study Categorizer")

# # === Load Documents from Azure ===
# @st.cache_resource
# def load_documents():
#     loader = AzureBlobStorageContainerLoader(
#         conn_str=os.getenv("AZURE_STORAGE_CONNECTION_STRING"),
#         container=os.getenv("AZURE_STORAGE_CONTAINER_NAME")
#     )
#     return loader.load()

# documents = load_documents()

# # === Prepare a list of supported file names ===
# supported_files = {}
# for doc in documents:
#     source_path = doc.metadata.get("source", "Unknown source")
#     file_name = os.path.basename(source_path)
#     if file_name.lower().endswith(('.pdf', '.pptx')):
#         supported_files[file_name] = doc.page_content

# # === Sidebar multiple file picker ===
# selected_files = st.sidebar.multiselect("📂 Select file(s)", options=list(supported_files.keys()))

# if selected_files:
#     for file in selected_files:
#         st.subheader(f"📝 Summary of `{file}`")
#         text = supported_files[file]

#         # 💡 Agent 1: Reader
#         summary = run_reader_agent(text)
#         st.markdown(summary)

#         # 📌 Agent 2: Categorizer
#         category = run_categorization_agent(summary, [])
#         st.markdown(f"**📁 Category:** {category}")

import streamlit as st
from dotenv import load_dotenv
import os
import fitz  # PyMuPDF
import docx
from pptx import Presentation
from langchain.text_splitter import RecursiveCharacterTextSplitter
from langchain.vectorstores import Chroma
from langchain_openai import AzureOpenAIEmbeddings
from langchain.chat_models import AzureChatOpenAI
from langchain.chains import RetrievalQA
from langchain.memory import ConversationBufferMemory

# 🌱 Load environment variables
load_dotenv()
api_key = os.getenv("AZURE_OPENAI_API_KEY")
api_base = os.getenv("AZURE_OPENAI_END_POINT")  # e.g. https://your-resource-name.openai.azure.com
api_version = os.getenv("API_VERSION")  # e.g. 2023-05-15
chat_deployment = os.getenv("MODEL_NAME")  # chat model name
embed_deployment = "text-embedding-ada-002"  # embedding model name

# 🎨 Streamlit UI setup
st.set_page_config(page_title="📚 Memory-Aware RAG Assistant", layout="wide")
st.title("📚 Memory-Aware RAG Q&A Assistant")

# 🚨 Validate environment
if not all([api_key, api_base, api_version, chat_deployment, embed_deployment]):
    st.error("❌ Missing one or more Azure OpenAI environment variables.")
    st.stop()

# 📄 File text extraction functions
@st.cache_data
def extract_text_from_pdf(pdf_file):
    doc = fitz.open(stream=pdf_file.read(), filetype="pdf")
    return "\n".join(page.get_text() for page in doc)

@st.cache_data
def extract_text_from_docx(docx_file):
    doc = docx.Document(docx_file)
    return "\n".join(para.text for para in doc.paragraphs)

@st.cache_data
def extract_text_from_txt(txt_file):
    return txt_file.read().decode("utf-8")

@st.cache_data
def extract_text_from_pptx(pptx_file):
    prs = Presentation(pptx_file)
    return "\n".join(shape.text for slide in prs.slides for shape in slide.shapes if hasattr(shape, "text"))

# 🧠 Create vector store
def create_vector_store(text):
    splitter = RecursiveCharacterTextSplitter(chunk_size=500, chunk_overlap=50)
    chunks = splitter.split_text(text)

    embedding_model = AzureOpenAIEmbeddings(
        azure_deployment=embed_deployment,
        openai_api_version=api_version,
        azure_endpoint=api_base,
        api_key=api_key
    )

    vector_store = Chroma.from_texts(
        chunks,
        embedding=embedding_model,
        collection_name="rag_store",
        persist_directory=".chromadb"
    )
    vector_store.persist()
    return vector_store

# 🔎 Set up QA chain
def get_qa_chain(vectorstore):
    memory = ConversationBufferMemory(
        memory_key="chat_history",
        return_messages=True,
        input_key="query",
        k=3
    )

    qa_chain = RetrievalQA.from_chain_type(
        llm=AzureChatOpenAI(
            openai_api_key=api_key,
            openai_api_base=api_base,
            openai_api_version=api_version,
            deployment_name=chat_deployment,
            temperature=0
        ),
        chain_type="stuff",
        retriever=vectorstore.as_retriever(),
        memory=memory,
        return_source_documents=False
    )
    return qa_chain

# 📤 File upload section
uploaded_files = st.file_uploader("📄 Upload Documents", type=["pdf", "docx", "txt", "pptx"], accept_multiple_files=True)

if uploaded_files:
    all_text = ""
    for file in uploaded_files:
        if file.type == "application/pdf":
            all_text += extract_text_from_pdf(file)
            st.success(f"✅ Processed PDF: {file.name}")
        elif file.type == "application/vnd.openxmlformats-officedocument.wordprocessingml.document":
            all_text += extract_text_from_docx(file)
            st.success(f"✅ Processed DOCX: {file.name}")
        elif file.type == "text/plain":
            all_text += extract_text_from_txt(file)
            st.success(f"✅ Processed TXT: {file.name}")
        elif file.type == "application/vnd.openxmlformats-officedocument.presentationml.presentation":
            all_text += extract_text_from_pptx(file)
            st.success(f"✅ Processed PPTX: {file.name}")
        else:
            st.warning(f"⚠️ Unsupported file type: {file.name}")

    # 💾 Initialize vector store and chain
    if "vectorstore" not in st.session_state:
        st.session_state.vectorstore = create_vector_store(all_text)
        st.session_state.qa_chain = get_qa_chain(st.session_state.vectorstore)
        st.session_state.chat_history = []

    # 💬 Sidebar chatbot
    st.sidebar.title("💬 Ask the Chatbot")
    user_question = st.sidebar.text_input("Ask a question:")

    if user_question:
        response = st.session_state.qa_chain({"query": user_question})
        answer = response.get('result', "❌ No answer found.")
        st.session_state.chat_history.append({"user": user_question, "bot": answer})
        st.markdown(f"**You asked:** {user_question}")
        st.markdown(f"**Answer:** {answer}")

    # 🧾 Chat history display
    if st.session_state.chat_history:
        with st.expander("📜 Chat History"):
            for i, chat in enumerate(st.session_state.chat_history, 1):
                st.markdown(f"**You {i}:** {chat['user']}")
                st.markdown(f"**Bot {i}:** {chat['bot']}")
else:
    st.info("👆 Upload one or more documents to begin.")

