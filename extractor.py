# extractor.py

import os
from PyPDF2 import PdfReader
from pptx import Presentation

def extract_text_pdf(file_path):
    reader = PdfReader(file_path)
    text = ""
    for page in reader.pages:
        curr = page.extract_text()
        if curr:
            text += curr.strip() + "\n"
    return text

def extract_text_pptx(file_path):
    prs = Presentation(file_path)
    text_runs = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text_runs.append(shape.text)
    return "\n".join(text_runs)

def extract_text(file_path):
    ext = file_path.lower().split('.')[-1]
    if ext == "pdf":
        return extract_text_pdf(file_path)
    elif ext == "pptx":
        return extract_text_pptx(file_path)
    else:
        raise ValueError("Unsupported file type: " + ext)
